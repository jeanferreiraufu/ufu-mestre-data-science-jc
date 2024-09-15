from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Trabalho de Programação Dinâmica"
subtitle_1.text = "Máxima Subsequência Crescente com a Soma Máxima (MSIS)\n\nAutor: Jean Carlo Alves Ferreira\nOrientadora: Marcia Aparecida Fernandes\nData: 11/06/2024"

# Slide 2: Introdução
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]
title_2.text = "Introdução"
content_2.text = "Definição do Problema: Encontrar a subsequência crescente de uma dada sequência de números cuja soma é maximizada.\n\nImportância: Aplicações práticas em análise de séries temporais, bioinformática e finanças."

# Slide 3: Conceitos Básicos
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]
title_3.text = "Conceitos Básicos"
content_3.text = "Programação Dinâmica: Técnica para resolver problemas complexos dividindo-os em subproblemas mais simples.\n\nSubsequência Crescente: Sequência de números em ordem crescente."

# Slide 4: Relação de Recorrência
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]
title_4.text = "Relação de Recorrência"
content_4.text = "Fórmula: MSIS(i) = max(MSIS(j) + arr[i]) para todo j < i e arr[j] < arr[i]\n\nObjetivo: Maximizar a soma da subsequência crescente."

# Slide 5: Exemplo
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]
title_5.text = "Exemplo"
content_5.text = "Entrada: [1, 101, 3, 2, 100, 4, 5]\n\nSaída: 106 (subsequência: [1, 3, 100])"

# Slide 6: Algoritmo Iterativo
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.placeholders[1]
title_6.text = "Algoritmo Iterativo"
content_6.text = "Descrição do Algoritmo:\n1. Inicializa o array MSIS.\n2. Atualiza MSIS para cada elemento.\n3. Calcula a soma máxima."

# Slide 7: Pseudocódigo Iterativo
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
content_7 = slide_7.placeholders[1]
title_7.text = "Pseudocódigo Iterativo"
content_7.text = """
function subSequenciaCrescenteSomaMaxima(arr):
    n = length(arr)
    if n == 0:
        return 0
    MSIS = array of size n
    for i = 0 to n-1:
        MSIS[i] = arr[i]
    for i = 1 to n-1:
        for j in range(i-1):
            if arr[i] > arr[j]:
                MSIS[i] = max(MSIS[i], MSIS[j] + arr[i])
    max_sum = MSIS[0]
    for i = 1 to n-1:
        if MSIS[i] > max_sum:
            max_sum = MSIS[i]
    return max_sum
"""

# Slide 8: Implementação em Java
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
content_8 = slide_8.placeholders[1]
title_8.text = "Implementação em Java"
content_8.text = """
public class SubSequenciaCrescenteSomaMaxima {
    static int maiorSomaE(int arr[]) {
        int n = arr.length;
        if (n==0) return 0;
        int MSIS[] = new int[n];
        for (int i = 0; i < n; i++) MSIS[i] = arr[i];
        for (int i = 1; i < n; i++) {
            for (int j = 0; j < i; j++) {
                if (arr[i] > arr[j] && MSIS[i] < MSIS[j] + arr[i]) {
                    MSIS[i] = MSIS[j] + arr[i];
                }
            }
        }
        int max = MSIS[0];
        for (int i = 1; i < n; i++) {
            if (MSIS[i] > max) max = MSIS[i];
        }
        return max;
    }
    public static void main(String[] args) {
        int arr[] = {1, 101, 3, 2, 100, 4, 5};
        System.out.println("A soma máxima da subsequência crescente é " + maiorSomaE(arr));
    }
}
"""

# Slide 9: Análise de Complexidade
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
content_9 = slide_9.placeholders[1]
title_9.text = "Análise de Complexidade"
content_9.text = "Complexidade de Tempo: O(n^2)\nJustificativa: Dois loops aninhados para atualizar o array MSIS."

# Slide 10: Algoritmo Recursivo
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title_10 = slide_10.shapes.title
content_10 = slide_10.placeholders[1]
title_10.text = "Algoritmo Recursivo"
content_10.text = "Descrição do Algoritmo:\n1. Calcula a soma máxima e a subsequência correspondente para uma posição i.\n2. Sem memorização, recalcula subproblemas várias vezes."

# Slide 11: Pseudocódigo Recursivo
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title_11 = slide_11.shapes.title
content_11 = slide_11.placeholders[1]
title_11.text = """
def subSequenciaCrescenteSomaMaxima(arr, i):
    if i == 0:
        return arr[0], [arr[0]]
    max_sum = arr[i]
    max_subsequence = [arr[i]]
    for j in range(i):
        prev_sum, prev_subsequence = subSequenciaCrescenteSomaMaxima(arr, j)
        if arr[i] > arr[j] and max_sum < prev_sum + arr[i]:
            max_sum = prev_sum + arr[i]
            max_subsequence = prev_subsequence + [arr[i]]
    return max_sum, max_subsequence
"""

# Slide 12: Comparativo de Execução
slide_12 = prs.slides.add_slide(prs.slide_layouts[1])
title_12 = slide_12.shapes.title
content_12 = slide_12.placeholders[1]
title_12.text = """
Tabela de Resultados:
Tamanho da Instância | Tempo Interativo (ms) | Tempo Recursivo (ms)
10000                | 155                   | 113
20000                | 542                   | 402
...                  | ...                   | ...
300000               | 93176                 | 93169
"""

# Slide 13: Conclusão
slide_13 = prs.slides.add_slide(prs.slide_layouts[1])
title_13 = slide_13.shapes.title
content_13 = slide_13.placeholders[1]
title_13.text = "Conclusão"
content_13.text = "Sumário: Abordagem clara e intuitiva para resolver o problema de MSIS usando programação dinâmica.\n\nRelevância: Aplicações práticas em várias áreas.\n\nComplexidade: Aceitável para muitos casos práticos."

# Slide 14: Referências
slide_14 = prs.slides.add_slide(prs.slide_layouts[1])
title_14 = slide_14.shapes.title
content_14 = slide_14.placeholders[1]
title_14.text = "Referências"
content_14.text = """
https://iq.opengenus.org/maximum-sum-increasing-subsequence/
https://ifacet.iitk.ac.in/knowledge-hub/data-structure-with-c/dynamic-programming-techniques-for-solving-optimization-problems/
https://www.educative.io/courses/grokking-dynamic-programming-a-deep-dive-using-java/maximum-sum-increasing-subsequence
"""

# Save the presentation
prs.save('presentation.pptx')
